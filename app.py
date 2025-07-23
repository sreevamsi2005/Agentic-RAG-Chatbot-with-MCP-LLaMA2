
import streamlit as st
import os
import uuid
import fitz  # PyMuPDF
import pandas as pd
from docx import Document
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import chromadb
from chromadb.utils import embedding_functions
from langchain.text_splitter import CharacterTextSplitter
import numpy as np
import requests

# ------------------- MCP Message Generator -------------------
def create_mcp_message(type, sender, receiver, trace_id, payload):
    return {
        "type": type,
        "sender": sender,
        "receiver": receiver,
        "trace_id": trace_id,
        "payload": payload
    }

# ------------------- Agent Classes -------------------
class IngestionAgent:
    model = SentenceTransformer('all-MiniLM-L6-v2')

    @staticmethod
    def parse_file(file_path):
        try:
            if file_path.endswith(".pdf"):
                text = ""
                with fitz.open(file_path) as doc:
                    for page in doc:
                        text += page.get_text()
                return text.strip()
            elif file_path.endswith(".csv"):
                df = pd.read_csv(file_path)
                return df.to_string(index=False)
            elif file_path.endswith(".docx"):
                doc = Document(file_path)
                return "\n".join([para.text for para in doc.paragraphs])
            elif file_path.endswith(".txt") or file_path.endswith(".md"):
                with open(file_path, "r", encoding="utf-8") as f:
                    return f.read()
            elif file_path.endswith(".pptx"):
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text.strip()
            else:
                return ""
        except Exception as e:
            return f"Error parsing file: {e}"

    @staticmethod
    def receive_message(message, router):
        file_path = message['payload'].get("file_path")
        if not file_path or not os.path.exists(file_path):
            return

        text = IngestionAgent.parse_file(file_path)
        splitter = CharacterTextSplitter(chunk_size=500, chunk_overlap=50)
        chunks = splitter.split_text(text)
        embeddings = IngestionAgent.model.encode(chunks).tolist()

        chroma_client = chromadb.Client()
        embedding_fn = embedding_functions.SentenceTransformerEmbeddingFunction(model_name="all-MiniLM-L6-v2")

        try:
            collection = chroma_client.get_or_create_collection(name="rag_collection", embedding_function=embedding_fn)
        except Exception as e:
            collection = chroma_client.create_collection(name="rag_collection", embedding_function=embedding_fn)

        ids = [str(uuid.uuid4()) for _ in range(len(chunks))]
        collection.add(documents=chunks, ids=ids)

        router.memory = {"chunks": chunks, "collection": collection}

        new_msg = create_mcp_message(
            type="DOC_PARSED",
            sender="IngestionAgent",
            receiver="RetrievalAgent",
            trace_id=message["trace_id"],
            payload={"query": message['payload'].get("query", "What are the KPIs?")}
        )
        router.send_msg(new_msg)


class RetrievalAgent:
    @staticmethod
    def receive_message(message, router):
        query = message['payload'].get("query", "What are the KPIs?")
        collection = router.memory["collection"]
        results = collection.query(query_texts=[query], n_results=3)
        top_chunks = results['documents'][0] if results['documents'] else []

        new_msg = create_mcp_message(
            type="RETRIEVAL_RESULT",
            sender="RetrievalAgent",
            receiver="LLMAgent",
            trace_id=message['trace_id'],
            payload={"retrieved_chunks": top_chunks, "query": query}
        )
        router.send_msg(new_msg)


class LLMAgent:
    @staticmethod
    def receive_message(message, router):
        chunks = message['payload'].get("retrieved_chunks", [])
        query = message['payload'].get("query", "")
        context_text = "\n".join(chunks)

        prompt = f"""
        You are a professional data analyst. Use the following context to answer the user's query.

        Context:
        {context_text}

        Question:
        {query}

        Provide a precise and clear answer below:
        """

        try:
            response = requests.post(
                "http://localhost:11434/api/generate",
                json={"model": "llama2:7b", "prompt": prompt, "stream": False}
            )
            data = response.json()
            answer = data.get("response", "[No response from LLM]")
        except Exception as e:
            answer = f"[Error communicating with Ollama] {e}"

        router.chat_history.append((query, answer))


# ------------------- MCP Router -------------------
class MCPRouter:
    def __init__(self):
        self.agent_registry = {
            "IngestionAgent": IngestionAgent,
            "RetrievalAgent": RetrievalAgent,
            "LLMAgent": LLMAgent,
        }
        self.memory = {}
        self.chat_history = []

    def send_msg(self, message):
        receiver = message.get("receiver")
        agent = self.agent_registry.get(receiver)
        if agent:
            agent.receive_message(message, self)


# ------------------- Streamlit UI -------------------
st.set_page_config(page_title="Agentic RAG Chatbot", layout="wide")
st.title("🤖 Agentic RAG Chatbot (MCP + RAG + LLaMA2)")

if "router" not in st.session_state:
    st.session_state.router = MCPRouter()

with st.sidebar:
    uploaded_file = st.file_uploader("📤 Upload a document", type=["pdf", "csv", "docx", "txt", "md", "pptx"])
    if uploaded_file:
        file_name = f"uploads/{uploaded_file.name}"
        os.makedirs("uploads", exist_ok=True)
        with open(file_name, "wb") as f:
            f.write(uploaded_file.read())
        st.success("✅ Uploaded: " + uploaded_file.name)

        if st.button("🚀 Ingest Document"):
            trace_id = str(uuid.uuid4())
            msg = create_mcp_message(
                type="UPLOAD_DOC",
                sender="UI",
                receiver="IngestionAgent",
                trace_id=trace_id,
                payload={"file_path": file_name, "query": "placeholder"}
            )
            st.session_state.router.send_msg(msg)
            st.success("🧠 Document processed and embedded.")

if st.session_state.router.memory:
    user_input = st.text_input("💬 Ask a question about the uploaded document:")
    if st.button("🔍 Submit Query") and user_input:
        trace_id = str(uuid.uuid4())
        msg = create_mcp_message(
            type="QUERY",
            sender="UI",
            receiver="RetrievalAgent",
            trace_id=trace_id,
            payload={"query": user_input}
        )
        st.session_state.router.send_msg(msg)

if st.session_state.router.chat_history:
    st.subheader("🧠 Chat History")
    for i, (q, a) in enumerate(st.session_state.router.chat_history[::-1]):
        st.markdown(f"**Q{i+1}:** {q}")
        st.markdown(f"**A{i+1}:** {a}")
